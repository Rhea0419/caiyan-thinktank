"""
采研智库——公众号跟踪 前端 v1.3
"""
import streamlit as st
import sqlite3, os, json, io, tempfile, re
from datetime import datetime, timedelta
import pandas as pd

st.set_page_config(page_title="采研智库·公众号跟踪", page_icon="📡", layout="wide")

DB_PATH = "/Users/rhea/Documents/Kapathy/commodities/data/articles.db"
ARCHIVE_DIR = "/Users/rhea/Documents/Kapathy/commodities/archive"
os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
os.makedirs(ARCHIVE_DIR, exist_ok=True)

def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS articles (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        url TEXT UNIQUE, title TEXT, source TEXT, author TEXT,
        pub_date TEXT, category TEXT, tags TEXT,
        summary TEXT, body TEXT, key_points TEXT,
        policy_regulation TEXT, industry_conference TEXT, key_data TEXT, core_argument TEXT,
        industry_impact TEXT, bid_impact TEXT, company_news TEXT, procurement_trend TEXT,
        created_at TEXT
    )''')
    c.execute('''CREATE TABLE IF NOT EXISTS archives (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        filename TEXT, upload_date TEXT, report_month TEXT,
        category TEXT, extracted_text TEXT, key_info TEXT, tags TEXT
    )''')
    for idx in ['idx_category','idx_date','idx_source']:
        try: c.execute(f'CREATE INDEX IF NOT EXISTS {idx} ON articles({idx.split("_")[1]})')
        except: pass
    conn.commit()
    conn.close()

init_db()

# === 文档解析 ===
def extract_text(uploaded_file):
    ext = os.path.splitext(uploaded_file.name)[1].lower()
    try:
        if ext in ['.txt','.md']:
            return uploaded_file.read().decode('utf-8', errors='ignore')[:10000]
        elif ext == '.docx':
            from docx import Document
            return '\n'.join(p.text for p in Document(uploaded_file).paragraphs)[:10000]
        elif ext == '.pdf':
            import pymupdf
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
                tmp.write(uploaded_file.read()); tmp.flush()
                doc = pymupdf.open(tmp.name)
                text = '\n'.join(page.get_text() for page in doc)
                doc.close(); os.unlink(tmp.name)
                return text[:10000]
        elif ext in ['.xlsx','.xls']:
            return pd.read_excel(uploaded_file).to_string(max_rows=50)[:10000]
        elif ext == '.pptx':
            from pptx import Presentation
            texts = []
            for slide in Presentation(uploaded_file).slides:
                for shape in slide.shapes:
                    if shape.has_text_frame: texts.append(shape.text_frame.text)
            return '\n'.join(texts)[:10000]
        return f"不支持: {ext}"
    except Exception as e:
        return f"解析失败: {str(e)}"

# === CSS ===
st.markdown("""
<style>
    [data-testid="stToolbar"], header[data-testid="stHeader"],
    [data-testid="stDecoration"], #MainMenu, footer,
    .stDeployButton { display: none !important; }
    .main-title { font-size: 24px; font-weight: 700; color: #1a1a1a; }
    .sub-title { font-size: 13px; color: #6B7280; margin-bottom: 16px; }
    .stat-card { background: white; border-radius: 10px; padding: 14px 18px; 
                 box-shadow: 0 1px 3px rgba(0,0,0,0.05); text-align: center; }
    .stat-num { font-size: 28px; font-weight: 700; color: #3B82F6; }
    .stat-label { font-size: 11px; color: #6B7280; }
    .article-card { background: white; border-radius: 8px; padding: 16px; 
                    margin: 8px 0; border-left: 3px solid #3B82F6;
                    box-shadow: 0 1px 2px rgba(0,0,0,0.04); }
    .article-card.风电 { border-left-color: #10B981; }
    .article-card.光伏 { border-left-color: #F59E0B; }
    .article-card.火电 { border-left-color: #EF4444; }
    .article-card.招投标 { border-left-color: #8B5CF6; }
    .article-card.selected { border: 2px solid #3B82F6; background: #EFF6FF; }
    .tag { display: inline-block; background: #EFF6FF; color: #3B82F6; 
           padding: 2px 8px; border-radius: 4px; font-size: 11px; margin: 2px; }
    .info-tag { display: inline-block; background: #FEF3C7; color: #92400E; 
                padding: 2px 8px; border-radius: 4px; font-size: 11px; margin: 2px; }
    .data-tag { display: inline-block; background: #ECFDF5; color: #065F46; 
                padding: 2px 8px; border-radius: 4px; font-size: 11px; margin: 2px; }
    .category-badge { display: inline-block; padding: 3px 10px; border-radius: 12px; 
                      font-size: 11px; font-weight: 600; color: white; }
    .category-badge.储能 { background: #3B82F6; }
    .category-badge.风电 { background: #10B981; }
    .category-badge.光伏 { background: #F59E0B; }
    .category-badge.火电 { background: #EF4444; }
    .category-badge.招投标 { background: #8B5CF6; }
    .source-item { font-size: 12px; padding: 3px 0; border-bottom: 1px solid #F3F4F6; }
    .source-item .name { font-weight: 600; }
    .source-item .date { color: #9CA3AF; font-size: 11px; }
    .section-label { font-weight: 700; margin-top: 6px; display: block; font-size: 12px; }
    .section-label.policy { color: #92400E; }
    .section-label.conference { color: #7C3AED; }
    .section-label.data { color: #065F46; }
    .section-label.argument { color: #9333EA; }
    .chat-msg { padding: 10px 14px; border-radius: 8px; margin: 6px 0; font-size: 13px; }
    .chat-msg.user { background: #EFF6FF; text-align: right; }
    .chat-msg.assistant { background: #F9FAFB; border: 1px solid #E5E7EB; }
    .chat-source { font-size: 10px; color: #9CA3AF; margin-top: 4px; }
</style>
""", unsafe_allow_html=True)

# === 导航 ===
if "view" not in st.session_state:
    st.session_state.view = "入库"

views = {"入库": "📥", "检索": "🔍", "月报": "📊", "归档": "🗄️", "导出": "📤"}
cols = st.columns([1,1,1,1,1,3])
for i, (k, v) in enumerate(views.items()):
    with cols[i]:
        if st.button(f"{v} {k}", key=f"nav_{k}", use_container_width=True,
                     type="primary" if st.session_state.view == k else "secondary"):
            st.session_state.view = k
            st.rerun()
st.markdown("---")

# ===================== 入库 =====================
if st.session_state.view == "入库":
    col_left, col_right = st.columns([3, 1])
    with col_left:
        st.markdown('<p class="main-title">📥 文章入库</p>', unsafe_allow_html=True)
        tab1, tab2 = st.tabs(["🔗 链接录入", "📄 文档上传"])
        with tab1:
            url_input = st.text_area("文章链接", height=80, placeholder="https://mp.weixin.qq.com/s/...")
            if st.button("🔍 提取文章信息", type="primary") and url_input:
                urls = [u.strip() for u in url_input.split('\n') if u.strip()]
                for url in urls:
                    conn = sqlite3.connect(DB_PATH); c = conn.cursor()
                    c.execute("SELECT id,title FROM articles WHERE url=?", (url,))
                    if c.fetchone(): st.warning(f"⏭️ 已存在: {url[:50]}..."); conn.close(); continue
                    conn.close()
                    st.session_state.setdefault('pending_urls',{})[url.split('/')[-1][:8]] = url
        with tab2:
            uploaded = st.file_uploader("上传文档", type=['docx','pdf','xlsx','xls','pptx','txt','md'])
            if uploaded:
                extracted = extract_text(uploaded)
                aid = str(hash(uploaded.name))[-8:]
                st.session_state.setdefault('pending_docs',{})[aid] = {'name':uploaded.name,'text':extracted}
                st.text_area("提取预览", extracted, height=200, key=f"pre_{aid}")
        
        pending = {**st.session_state.get('pending_urls',{}), **st.session_state.get('pending_docs',{})}
        for aid, data in pending.items():
            url = data if isinstance(data, str) else ''
            dn = data.get('name','') if isinstance(data, dict) else ''
            dt = data.get('text','') if isinstance(data, dict) else ''
            label = f"📄 {dn}" if dn else f"📄 {url[:60]}..."
            with st.expander(label, expanded=True):
                c1, c2 = st.columns(2)
                with c1: title = st.text_input("标题", key=f"t_{aid}"); source = st.text_input("公众号", key=f"s_{aid}")
                with c2: pub_date = st.date_input("日期", key=f"d_{aid}")
                category = st.selectbox("行业", ["储能","风电","光伏","火电","招投标","其他"], key=f"c_{aid}")
                tags = st.text_input("标签", key=f"tg_{aid}")
                summary = st.text_area("摘要（100字）", height=50, key=f"sm_{aid}", value=dt[:100] if dt else '')
                st.markdown("#### 🔬 深度分析")
                a1, a2 = st.columns(2)
                with a1: pr = st.text_area("政策法规", height=45, key=f"pr_{aid}"); kd = st.text_area("关键数据", height=45, key=f"kd_{aid}")
                with a2: ic = st.text_area("行业会议", height=45, key=f"ic_{aid}"); ca = st.text_area("核心论点", height=45, key=f"ca_{aid}")
                with st.expander("📊 月报分析"):
                    b1, b2 = st.columns(2)
                    with b1: ii = st.text_area("行业影响", height=45, key=f"ii_{aid}"); cn = st.text_area("企业动态", height=45, key=f"cn_{aid}")
                    with b2: bi = st.text_area("招投标影响", height=45, key=f"bi_{aid}"); pt = st.text_area("采购趋势", height=45, key=f"pt_{aid}")
                if st.button("💾 保存", key=f"save_{aid}"):
                    conn = sqlite3.connect(DB_PATH); c = conn.cursor()
                    c.execute('''INSERT OR REPLACE INTO articles VALUES (NULL,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)''',
                        (url or f"file://{dn}",title,source,'',str(pub_date),category,tags,summary,dt,'',pr,ic,kd,ca,ii,bi,cn,pt,datetime.now().isoformat()))
                    conn.commit(); conn.close()
                    st.success(f"✅ {title}"); st.rerun()
    with col_right:
        st.markdown("### 📋 已导入清单")
        conn = sqlite3.connect(DB_PATH); c = conn.cursor()
        c.execute("SELECT source, COUNT(*), MAX(pub_date) FROM articles WHERE source!='' GROUP BY source ORDER BY MAX(pub_date) DESC")
        for src,cnt,latest in c.fetchall():
            st.markdown(f"""<div class="source-item"><span class="name">{src}</span>({cnt}篇)<br><span class="date">最新:{latest}</span></div>""", unsafe_allow_html=True)
        st.metric("总文章", sum(r[1] for r in c.fetchall() if False) or c.execute("SELECT COUNT(*) FROM articles").fetchone()[0])
        conn.close()

# ===================== 检索 =====================
elif st.session_state.view == "检索":
    st.markdown('<p class="main-title">🔍 文章检索</p>', unsafe_allow_html=True)
    conn = sqlite3.connect(DB_PATH); c = conn.cursor()
    c.execute("SELECT COUNT(*), COUNT(DISTINCT source), COUNT(DISTINCT category) FROM articles")
    total, sources, cats = c.fetchone()
    sts = st.columns(4)
    for i, (n, v) in enumerate([("文章总数",total),("公众号数",sources),("行业分类",cats),("本周新增",c.execute("SELECT COUNT(*) FROM articles WHERE pub_date>=date('now','-7 days')").fetchone()[0])]):
        with sts[i]: st.markdown(f'<div class="stat-card"><div class="stat-num">{v}</div><div class="stat-label">{n}</div></div>', unsafe_allow_html=True)
    conn.close(); st.markdown("---")
    
    f1,f2,f3,f4,f5 = st.columns([2,1.5,1,1,0.8])
    with f1: search = st.text_input("🔎", placeholder="关键词搜索...", label_visibility="collapsed")
    with f2: info_types = st.multiselect("类别", ["政策法规","行业会议","关键数据","核心论点"], default=[], placeholder="信息类别", label_visibility="collapsed")
    with f3: date_from = st.date_input("从", value=None, label_visibility="collapsed")
    with f4: date_to = st.date_input("至", value=None, label_visibility="collapsed")
    with f5: st.button("🔍", use_container_width=True)
    
    f6,f7 = st.columns(2)
    with f6:
        conn = sqlite3.connect(DB_PATH); c = conn.cursor()
        c.execute("SELECT DISTINCT category FROM articles WHERE category!=''")
        cat_filter = st.multiselect("行业", [r[0] for r in c.fetchall()], default=[], placeholder="行业分类", label_visibility="collapsed")
        conn.close()
    with f7:
        conn = sqlite3.connect(DB_PATH); c = conn.cursor()
        c.execute("SELECT DISTINCT source FROM articles WHERE source!=''")
        src_filter = st.multiselect("公众号", [r[0] for r in c.fetchall()], default=[], placeholder="公众号", label_visibility="collapsed")
        conn.close()
    
    conn = sqlite3.connect(DB_PATH)
    query = "SELECT * FROM articles WHERE 1=1"; params = []
    if search: query += " AND (title||summary||tags||policy_regulation||industry_conference||key_data||core_argument LIKE ?)"; params.append(f"%{search}%")
    info_map = {"政策法规":"policy_regulation","行业会议":"industry_conference","关键数据":"key_data","核心论点":"core_argument"}
    for it in info_types: query += f" AND {info_map[it]} IS NOT NULL AND {info_map[it]}!='' AND {info_map[it]}!='None'"
    if cat_filter: query += " AND category IN ("+",".join(["?"]*len(cat_filter))+")"; params.extend(cat_filter)
    if src_filter: query += " AND source IN ("+",".join(["?"]*len(src_filter))+")"; params.extend(src_filter)
    if date_from: query += " AND pub_date>=?"; params.append(str(date_from))
    if date_to: query += " AND pub_date<=?"; params.append(str(date_to))
    query += " ORDER BY pub_date DESC LIMIT 50"
    df = pd.read_sql(query, conn, params=params); conn.close()
    st.caption(f"找到 {len(df)} 篇")
    
    if df.empty: st.info("📭 无匹配文章")
    else:
        for _, row in df.iterrows():
            cat = row.get('category','储能')
            tags_html = ''.join(f'<span class="tag">{t.strip()}</span>' for t in (row.get('tags','') or '').split(',') if t.strip())
            sections = []
            for field, label, cls in [('policy_regulation','📜 政策','policy'),('industry_conference','📅 会议','conference'),('key_data','📊 数据','data'),('core_argument','💡 论点','argument')]:
                v = str(row.get(field,''))
                if v and v != 'None': sections.append(f'<span class="section-label {cls}">{label}</span><small>{v[:120]}</small>')
            url_link = row.get('url','')
            link_line = f'📎 <a href="{url_link}" target="_blank">原文</a>' if url_link else ''
            sections_joined = '<br>'.join(sections)
            st.markdown(f"""<div class="article-card {cat}"><span class="category-badge {cat}">{cat}</span> <strong>{row['title']}</strong><br>
                <small>📅 {row.get('pub_date','?')} | 📡 {row.get('source','?')} | {link_line}</small><br>
                <small>📝 摘要：{row.get('summary','')[:100]}</small><br>{tags_html}{sections_joined}</div>""", unsafe_allow_html=True)
            st.markdown("---")

# ===================== 月报（AI问答 + PDF导出）=====================
elif st.session_state.view == "月报":
    tab_m1, tab_m2, tab_m3 = st.tabs(["📊 月报生成", "🤖 AI问答", "📄 PDF排版"])
    
    # --- Tab 1: 月报生成 ---
    with tab_m1:
        conn = sqlite3.connect(DB_PATH); c = conn.cursor()
        c.execute("SELECT DISTINCT substr(pub_date,1,7) FROM articles WHERE pub_date!='' ORDER BY 1 DESC")
        months = [r[0] for r in c.fetchall()] or [datetime.now().strftime("%Y-%m")]
        selected_month = st.selectbox("选择月份", months, index=0)
        c.execute("SELECT * FROM articles WHERE substr(pub_date,1,7)=? ORDER BY pub_date DESC", (selected_month,))
        articles = c.fetchall(); conn.close()
        if not articles: st.warning(f"📭 {selected_month} 暂无文章")
        else:
            cats_count = {}
            for a in articles:
                cat = a[6] or '其他'
                cats_count[cat] = cats_count.get(cat, 0) + 1
            cols = st.columns(5)
            for i, cn in enumerate(["储能","风电","光伏","火电","招投标"]):
                with cols[i]: st.markdown(f'<div class="stat-card"><div class="stat-num">{cats_count.get(cn,0)}</div><div class="stat-label">{cn}</div></div>', unsafe_allow_html=True)
            st.markdown("---")
            report = [f"# 采购月报 —— {selected_month}\n", f"> 📅 {datetime.now().strftime('%Y-%m-%d')} | 📊 {len(articles)}篇\n", "---\n## 📈 行业动态\n"]
            for cat in ["储能","风电","光伏","火电","招投标"]:
                ca = [a for a in articles if a[6]==cat]
                if ca:
                    report.append(f"### {cat}（{len(ca)}篇）\n")
                    for a in ca[:5]:
                        report.append(f"- **{a[2]}** | {a[3]} | {a[8][:60] if a[8] else ''}\n")
                        if len(a)>18 and a[18]: report.append(f"  - 📊 {a[18][:100]}\n")
                    report.append("")
            report_text = '\n'.join(report)
            st.markdown(report_text)
            st.download_button("📥 下载 Markdown", report_text, f"采购月报-{selected_month}.md")
    
    # --- Tab 2: AI问答 ---
    with tab_m2:
        st.markdown("### 🤖 AI 问答（基于收录文章）")
        st.caption("在已收录的文章库中检索答案，不引用外部信息")
        
        # 信息源筛选
        with st.expander("📌 信息源筛选", expanded=True):
            sf1, sf2, sf3 = st.columns(3)
            with sf1:
                conn = sqlite3.connect(DB_PATH); c = conn.cursor()
                c.execute("SELECT DISTINCT substr(pub_date,1,7) FROM articles WHERE pub_date!='' ORDER BY 1 DESC")
                ai_month = st.selectbox("月份", ["全部"]+[r[0] for r in c.fetchall()], key="ai_m")
                conn.close()
            with sf2:
                conn = sqlite3.connect(DB_PATH); c = conn.cursor()
                c.execute("SELECT DISTINCT category FROM articles WHERE category!=''")
                ai_cat = st.multiselect("行业", [r[0] for r in c.fetchall()], key="ai_cat")
                conn.close()
            with sf3:
                conn = sqlite3.connect(DB_PATH); c = conn.cursor()
                c.execute("SELECT DISTINCT source FROM articles WHERE source!=''")
                ai_src = st.multiselect("公众号", [r[0] for r in c.fetchall()], key="ai_src")
                conn.close()
            
            # 勾选特定文章
            conn = sqlite3.connect(DB_PATH)
            q = "SELECT id, title, source, pub_date FROM articles WHERE 1=1"; params = []
            if ai_month != "全部": q += " AND substr(pub_date,1,7)=?"; params.append(ai_month)
            if ai_cat: q += " AND category IN ("+",".join(["?"]*len(ai_cat))+")"; params.extend(ai_cat)
            if ai_src: q += " AND source IN ("+",".join(["?"]*len(ai_src))+")"; params.extend(ai_src)
            q += " ORDER BY pub_date DESC LIMIT 30"
            arts = pd.read_sql(q, conn, params=params); conn.close()
            
            if not arts.empty:
                selected_ids = st.multiselect("或勾选特定文章", 
                    [f"{r['id']}: {r['title'][:40]} ({r['source']})" for _, r in arts.iterrows()],
                    key="ai_sel")
        
        # 聊天区
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []
        
        for msg in st.session_state.chat_history:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])
                if msg.get("sources"):
                    st.caption(f"📚 来源: {msg['sources']}")
        
        question = st.chat_input("基于收录文章提问...")
        if question:
            st.session_state.chat_history.append({"role":"user","content":question})
            
            # 检索相关文章
            conn = sqlite3.connect(DB_PATH)
            q2 = "SELECT * FROM articles WHERE 1=1"; params2 = []
            if selected_ids:
                ids = [int(s.split(':')[0]) for s in selected_ids]
                q2 += " AND id IN ("+",".join(["?"]*len(ids))+")"; params2.extend(ids)
            else:
                if ai_month!="全部": q2 += " AND substr(pub_date,1,7)=?"; params2.append(ai_month)
                if ai_cat: q2 += " AND category IN ("+",".join(["?"]*len(ai_cat))+")"; params2.extend(ai_cat)
                if ai_src: q2 += " AND source IN ("+",".join(["?"]*len(ai_src))+")"; params2.extend(ai_src)
                keywords = question.replace('？','').replace('?','').replace('，',' ').replace('、',' ').split()
                for kw in keywords[:5]:
                    if len(kw) >= 2:
                        q2 += " AND (title LIKE ? OR summary LIKE ? OR key_data LIKE ? OR core_argument LIKE ? OR policy_regulation LIKE ?)"
                        params2.extend([f"%{kw}%"]*5)
            q2 += " LIMIT 8"
            
            df2 = pd.read_sql(q2, conn, params=params2); conn.close()
            
            if df2.empty:
                answer = "📭 在选定的信息源中未找到相关内容。请扩大筛选范围或更换关键词。"
                sources = ""
            else:
                # 构建文章上下文
                context_parts = []
                src_names = []
                for _, r in df2.iterrows():
                    src_names.append(f"《{r['title'][:40]}》({r['source']},{r['pub_date']})")
                    ctx = f"【文章】{r['title']}\n来源：{r['source']} | 日期：{r['pub_date']} | 分类：{r['category']}\n摘要：{r.get('summary','')}"
                    for field, label in [('key_data','关键数据'),('core_argument','核心论点'),('policy_regulation','政策法规'),('industry_conference','行业会议')]:
                        v = str(r.get(field,''))
                        if v and v != 'None': ctx += f"\n{label}：{v}"
                    context_parts.append(ctx)
                
                context = "\n\n---\n\n".join(context_parts)
                
                # 调用 LLM
                try:
                    import urllib.request, ssl
                    ssl._create_default_https_context = ssl._create_unverified_context
                    
                    # 读取 API key
                    with open(os.path.expanduser("~/.hermes/.env")) as f:
                        env_vars = dict(line.strip().split("=", 1) for line in f if "=" in line and not line.startswith("#"))
                    api_key = env_vars.get("ANTHROPIC_API_KEY","")
                    
                    system_prompt = "你是采研智库的AI助手，基于已收录的新能源行业文章回答问题。只使用提供的文章内容作答，不要引用外部知识。如果文章信息不足以回答问题，请如实说明。用中文回答，简洁专业。"
                    user_prompt = f"请基于以下收录文章回答用户问题。\n\n=== 相关文章 ===\n{context}\n\n=== 用户问题 ===\n{question}\n\n请综合以上文章信息，给出准确、有据可循的回答，标注引用来源。"
                    
                    req_body = json.dumps({
                        "model": "deepseek-chat",
                        "max_tokens": 1500,
                        "temperature": 0.3,
                        "messages": [
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_prompt}
                        ]
                    }).encode()
                    
                    req = urllib.request.Request("https://api.deepseek.com/v1/chat/completions",
                        data=req_body,
                        headers={"Content-Type":"application/json", "Authorization": f"Bearer {api_key}"})
                    resp = urllib.request.urlopen(req, timeout=30)
                    resp_data = json.loads(resp.read())
                    answer = resp_data["choices"][0]["message"]["content"]
                except Exception as e:
                    # 降级：使用关键词匹配
                    parts = []
                    for _, r in df2.iterrows():
                        parts.append(f"**{r['title']}** ({r['source']}, {r['pub_date']})\n{r.get('summary','')[:150]}")
                        kd = str(r.get('key_data',''))
                        if kd and kd != 'None': parts.append(f"📊 {kd[:150]}")
                    answer = f"⚠️ AI 接口暂时不可用（{str(e)[:80]}），以下为关键词匹配结果：\n\n" + "\n\n".join(f"---\n{p}" for p in parts[:5])
                
                sources = "、".join(src_names[:5])
            
            st.session_state.chat_history.append({"role":"assistant","content":answer,"sources":sources})
            st.rerun()
    
    # --- Tab 3: PDF排版 ---
    with tab_m3:
        st.markdown("### 📄 PDF 排版导出")
        st.caption("期刊风格排版，生成专业月报PDF")
        
        conn = sqlite3.connect(DB_PATH); c = conn.cursor()
        c.execute("SELECT DISTINCT substr(pub_date,1,7) FROM articles WHERE pub_date!='' ORDER BY 1 DESC")
        pdf_month = st.selectbox("选择月份", [r[0] for r in c.fetchall()] or [datetime.now().strftime("%Y-%m")], key="pdf_m")
        conn.close()
        
        col1, col2 = st.columns(2)
        with col1:
            report_title = st.text_input("报告标题", f"采购月报 —— {pdf_month}")
        with col2:
            subtitle = st.text_input("副标题", "采研智库 · 采购研究中心")
        
        if st.button("📥 生成PDF", type="primary"):
            try:
                from reportlab.lib.pagesizes import A4
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.lib.units import mm
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
                from reportlab.lib import colors
                from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
                from reportlab.pdfbase import pdfmetrics
                from reportlab.pdfbase.ttfonts import TTFont
                
                # 注册中文字体
                font_path = '/System/Library/Fonts/STHeiti Medium.ttc'
                pdfmetrics.registerFont(TTFont('ZhFont', font_path))
                
                output_path = os.path.join(ARCHIVE_DIR, f"采购月报-{pdf_month}.pdf")
                doc = SimpleDocTemplate(output_path, pagesize=A4,
                    leftMargin=20*mm, rightMargin=20*mm, topMargin=25*mm, bottomMargin=20*mm)
                
                styles = getSampleStyleSheet()
                zh_normal = ParagraphStyle('ZhNormal', parent=styles['Normal'], fontName='ZhFont', fontSize=10, leading=16, alignment=TA_JUSTIFY)
                zh_title = ParagraphStyle('ZhTitle', parent=styles['Title'], fontName='ZhFont', fontSize=22, leading=28, alignment=TA_CENTER, spaceAfter=12)
                zh_h2 = ParagraphStyle('ZhH2', parent=styles['Heading2'], fontName='ZhFont', fontSize=14, leading=20, spaceBefore=16, spaceAfter=8)
                zh_h3 = ParagraphStyle('ZhH3', parent=styles['Heading3'], fontName='ZhFont', fontSize=12, leading=16, spaceBefore=12, spaceAfter=6)
                zh_small = ParagraphStyle('ZhSmall', parent=styles['Normal'], fontName='ZhFont', fontSize=8, leading=12, textColor=colors.HexColor('#6B7280'))
                
                story = []
                # 标题
                story.append(Paragraph(report_title, zh_title))
                story.append(Paragraph(subtitle, ParagraphStyle('Sub', parent=zh_small, alignment=TA_CENTER, fontSize=10)))
                story.append(Paragraph(f"生成日期: {datetime.now().strftime('%Y-%m-%d')}", zh_small))
                story.append(Spacer(1, 10*mm))
                story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor('#3B82F6')))
                story.append(Spacer(1, 8*mm))
                
                # 内容
                conn = sqlite3.connect(DB_PATH)
                df_m = pd.read_sql("SELECT * FROM articles WHERE substr(pub_date,1,7)=? ORDER BY pub_date DESC", conn, params=(pdf_month,))
                conn.close()
                
                story.append(Paragraph(f"📊 本期收录文章 {len(df_m)} 篇", zh_h2))
                
                for cat in ["储能","风电","光伏","火电","招投标"]:
                    cat_df = df_m[df_m['category'] == cat] if not df_m.empty else pd.DataFrame()
                    if not cat_df.empty:
                        story.append(Paragraph(f"■ {cat}（{len(cat_df)}篇）", zh_h3))
                        for _, r in cat_df.head(3).iterrows():
                            story.append(Paragraph(f"<b>{r['title']}</b> | {r['source']} | {r['pub_date']}", zh_normal))
                            if r.get('summary'): story.append(Paragraph(r['summary'][:200], zh_small))
                            if r.get('key_data') and str(r['key_data']) != 'None':
                                story.append(Paragraph(f"📊 {r['key_data'][:200]}", zh_small))
                            story.append(Spacer(1, 3*mm))
                doc.build(story)
                st.success(f"✅ PDF已生成: {output_path}")
                with open(output_path, 'rb') as f: st.download_button("📥 下载PDF", f, f"采购月报-{pdf_month}.pdf")
            except Exception as e:
                st.error(f"PDF生成失败: {e}")

# ===================== 归档 =====================
elif st.session_state.view == "归档":
    tab_a1, tab_a2 = st.tabs(["📤 上传归档", "🔍 归档检索"])
    
    with tab_a1:
        st.markdown("### 📤 往期月报归档")
        archive_file = st.file_uploader("上传月报PDF", type=['pdf'], key="arch_up")
        if archive_file:
            arch_month = st.text_input("报告月份", value=datetime.now().strftime("%Y-%m"))
            arch_text = extract_text(archive_file)
            st.text_area("提取内容", arch_text[:2000], height=200)
            arch_tags = st.text_input("标签")
            if st.button("📥 归档保存"):
                # 保存文件
                fname = f"月报-{arch_month}-{archive_file.name}"
                fpath = os.path.join(ARCHIVE_DIR, fname)
                with open(fpath, 'wb') as f: f.write(archive_file.getvalue())
                conn = sqlite3.connect(DB_PATH)
                conn.execute("INSERT INTO archives (filename,upload_date,report_month,category,extracted_text,key_info,tags) VALUES (?,?,?,?,?,?,?)",
                    (fname, datetime.now().isoformat(), arch_month, '', arch_text[:5000], '', arch_tags))
                conn.commit(); conn.close()
                st.success(f"✅ 已归档: {fname}")
    
    with tab_a2:
        st.markdown("### 🔍 归档检索")
        conn = sqlite3.connect(DB_PATH)
        search_a = st.text_input("搜索归档", placeholder="关键词...", key="arch_search")
        query_a = "SELECT * FROM archives WHERE 1=1"; params_a = []
        if search_a: query_a += " AND (filename LIKE ? OR extracted_text LIKE ? OR tags LIKE ? OR report_month LIKE ?)"; params_a.extend([f"%{search_a}%"]*4)
        query_a += " ORDER BY upload_date DESC LIMIT 30"
        df_a = pd.read_sql(query_a, conn, params=params_a); conn.close()
        st.caption(f"找到 {len(df_a)} 份归档")
        for _, r in df_a.iterrows():
            with st.expander(f"📄 {r['filename']} ({r['report_month']})"):
                st.caption(f"上传: {r['upload_date']} | 标签: {r.get('tags','')}")
                st.text(r.get('extracted_text','')[:1000])
                fpath = os.path.join(ARCHIVE_DIR, r['filename'])
                if os.path.exists(fpath):
                    with open(fpath, 'rb') as f: st.download_button("📥 下载", f, r['filename'], key=f"dl_{r['id']}")

# ===================== 导出 =====================
elif st.session_state.view == "导出":
    st.markdown('<p class="main-title">📤 数据导出</p>', unsafe_allow_html=True)
    conn = sqlite3.connect(DB_PATH)
    df = pd.read_sql("SELECT title,source,pub_date,category,tags,summary,policy_regulation,industry_conference,key_data,core_argument,url FROM articles ORDER BY pub_date DESC", conn)
    conn.close()
    if df.empty: st.info("暂无数据")
    else:
        st.dataframe(df, use_container_width=True, hide_index=True)
        c1, c2, c3 = st.columns(3)
        with c1: st.download_button("📥 CSV", df.to_csv(index=False), "articles.csv")
        with c2: st.download_button("📥 JSON", df.to_json(orient="records", force_ascii=False), "articles.json")
        with c3: st.metric("总文章数", len(df))

st.caption("采研智库 v1.3 | 采购研究中心")
